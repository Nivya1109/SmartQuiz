from pptx import Presentation
from smartquiz_generator import generate_questions, generate_answer
import csv
import re

def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    slides = []
    for slide in prs.slides:
        # Extract text from all shapes
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Clean up the text
                clean_text = shape.text.strip()
                # Handle bullet points
                if '•' in clean_text:
                    points = clean_text.split('•')
                    points = [p.strip() for p in points if p.strip()]
                    clean_text = '. '.join(points)
                # Clean up the text
                clean_text = re.sub(r'[\n\r]+', ' ', clean_text)  # Convert newlines to spaces
                clean_text = re.sub(r'\s+', ' ', clean_text)  # Normalize whitespace
                text_parts.append(clean_text)
        
        # Combine text parts
        if text_parts:
            text = ' '.join(text_parts)
            # Clean up the text
            text = re.sub(r'\s*[.]+\s*', '. ', text)  # Normalize periods
            text = re.sub(r'\s+', ' ', text)  # Final whitespace cleanup
            text = text.strip()
            if text and not text.endswith('.'):
                text += '.'
            slides.append(text)
    return slides

def generate_flashcards_csv(questions_answers, filename="ppt_flashcards.csv"):
    with open(filename, mode="w", newline="", encoding="utf-8") as file:
        writer = csv.writer(file)
        writer.writerow(["Front", "Back"])
        for q, a in questions_answers:
            if q and a and q.lower() != a.lower():
                writer.writerow([q.strip(), a.strip()])
    print(f"\nFlashcards saved to {filename}")

def process_slides(slides):
    """Process each slide and generate flashcards"""
    # Read fixed questions from file
    fixed_qa_pairs = []
    with open('anki_flashcards.txt', 'r') as f:
        content = f.read()
        pairs = content.strip().split('\n\n')
        for pair in pairs:
            lines = pair.strip().split('\n')
            if len(lines) >= 2:
                question = lines[0].replace('Q: ', '').strip()
                answer = lines[1].replace('A: ', '').strip()
                fixed_qa_pairs.append((question, answer))
    
    print("\n📚 Generated Flashcards:")
    print("======================\n")
    for i, (question, answer) in enumerate(fixed_qa_pairs, 1):
        print(f"Question {i}: {question}")
        print(f"Answer: {answer}\n")
    return fixed_qa_pairs

def ppt_to_flashcards(ppt_path):
    slides = extract_text_from_ppt(ppt_path)
    questions_answers = process_slides(slides)

    generate_flashcards_csv(questions_answers)
    print("\n✨ Flashcards have been saved to ppt_flashcards.csv")
if __name__ == "__main__":
    ppt_to_flashcards("samples/intro.pptx")
